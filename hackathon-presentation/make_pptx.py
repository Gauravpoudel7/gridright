"""
GridRight hackathon presentation -> GridRight.pptx
Seller + Operator model (NOT peer-to-peer).
Run: python make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ────────────────────────────────────────────────────
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
GREEN_DARK  = RGBColor(0x15, 0x80, 0x3d)
GREEN_SOFT  = RGBColor(0xec, 0xfd, 0xf5)
GREEN_BORD  = RGBColor(0xbb, 0xf7, 0xd0)
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)
AMBER_SOFT  = RGBColor(0xfe, 0xf3, 0xc7)
AMBER_BORD  = RGBColor(0xfc, 0xd3, 0x4d)
AMBER_INK   = RGBColor(0x78, 0x35, 0x0f)
PURPLE      = RGBColor(0x7c, 0x3a, 0xed)
PURPLE_SOFT = RGBColor(0xed, 0xe9, 0xfe)
PURPLE_BORD = RGBColor(0xc4, 0xb5, 0xfd)
INK         = RGBColor(0x0f, 0x17, 0x2a)
MUTED       = RGBColor(0x47, 0x55, 0x69)
LINE        = RGBColor(0xe2, 0xe8, 0xf0)
WHITE       = RGBColor(0xff, 0xff, 0xff)
SLATE_SOFT  = RGBColor(0xf8, 0xfa, 0xfc)
TITLE_BG    = RGBColor(0xf0, 0xfd, 0xf4)

FONT = "Segoe UI"

# ── Deck ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════
# CORE HELPERS
# ══════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _kill_shadow(shape):
    """python-pptx autoshapes inherit an ugly default shadow; strip it."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1.25), rounded=False):
    kind = 5 if rounded else 1          # ROUNDED_RECTANGLE / RECTANGLE
    s = slide.shapes.add_shape(kind, l, t, w, h)
    _kill_shadow(s)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def left_bar(slide):
    half = H // 2
    rect(slide, 0, 0,    Inches(0.16), half,   fill=GREEN)
    rect(slide, 0, half, Inches(0.16), H-half, fill=AMBER)

def _style_run(run, size, bold, color, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT

def txt(slide, text, l, t, w, h, size=24, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False, anchor=None, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        _style_run(r, size, bold, color, italic)
    return tb

def eyebrow(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.32), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    _style_run(r, 13, True, GREEN)
    r._r.get_or_add_rPr().set('spc', '1900')      # letter-spacing
    return tb

def heading(slide, text, t=Inches(0.78), size=42, color=INK,
            l=Inches(0.5), w=Inches(12.3)):
    tb = slide.shapes.add_textbox(l, t, w, Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        _style_run(r, size, True, color)
    return tb

def bullets(slide, items, l, t, w, h, size=21, color=INK,
            check=True, gap=Pt(7)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        pr = p._p.get_or_add_pPr()
        pr.set('marL', str(int(Pt(30)))); pr.set('indent', str(-int(Pt(30))))
        r1 = p.add_run(); r1.text = ("✓   " if check else "•   ")
        _style_run(r1, size, True, GREEN if check else MUTED)
        r2 = p.add_run(); r2.text = item
        _style_run(r2, size, False, color)
    return tb

def pill(slide, text, l, t, fill=GREEN, text_color=WHITE, size=15):
    """Auto-width rounded pill. Returns x for the next pill."""
    w = Inches(0.0092 * size * len(text) + 0.55)
    h = Inches(0.028 * size + 0.10)
    s = rect(slide, l, t, w, h, fill=fill, rounded=True)
    s.adjustments[0] = 0.5
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _style_run(r, size, True, text_color)
    return l + w + Inches(0.2)

def card(slide, l, t, w, h, fill=GREEN_SOFT, border=GREEN_BORD):
    return rect(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(1.5))

def circle_num(slide, label, l, t, d=Inches(0.55), fill=GREEN, size=20):
    s = slide.shapes.add_shape(9, l, t, d, d)     # OVAL
    _kill_shadow(s)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    _style_run(r, size, True, WHITE)
    return s

def callout(slide, text, l, t, w, h, fill=AMBER_SOFT, bar=AMBER, size=19):
    rect(slide, l, t, w, h, fill=fill)           # body first
    rect(slide, l, t, Inches(0.11), h, fill=bar) # accent bar ON TOP (fixed bug)
    tb = slide.shapes.add_textbox(l + Inches(0.3), t, w - Inches(0.5), h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        _style_run(r, size, False, INK)
    return tb

def footer(slide, n, total=12):
    txt(slide, f"{n} / {total}", Inches(11.9), Inches(7.06),
        Inches(1.2), Inches(0.35), size=12, color=MUTED, align=PP_ALIGN.RIGHT)

def flow_node(slide, icon, label, sub, l, t, w, h, fill, border):
    rect(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(1.5))
    txt(slide, icon, l, t + Inches(0.16), w, Inches(0.55),
        size=26, align=PP_ALIGN.CENTER)
    txt(slide, label, l + Inches(0.06), t + Inches(0.72), w - Inches(0.12),
        Inches(0.55), size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(slide, sub, l + Inches(0.08), t + Inches(1.28), w - Inches(0.16),
        Inches(0.6), size=12, color=MUTED, align=PP_ALIGN.CENTER)

def arrow(slide, l, t):
    txt(slide, "→", l, t, Inches(0.4), Inches(0.4),
        size=24, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, TITLE_BG); left_bar(s)
eyebrow(s, "Hackathon 2026  ·  Community Energy")
heading(s, "Your neighbor has solar.\nYou're paying full price.", t=Inches(0.85), size=46)
txt(s, "What if you could share it — fairly, transparently, on-chain?",
    Inches(0.5), Inches(2.95), Inches(12), Inches(0.6),
    size=25, bold=True, color=GREEN_DARK)

rect(s, Inches(0.5), Inches(3.65), Inches(3.75), Inches(0.7), fill=GREEN, rounded=True)
txt(s, "Introducing  GridRight", Inches(0.5), Inches(3.65), Inches(3.75), Inches(0.7),
    size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

txt(s, "Community solar pool  ·  Seller contributions  ·  Operator-governed  ·  Solana settlement",
    Inches(0.5), Inches(4.55), Inches(12.5), Inches(0.5),
    size=17, color=MUTED, italic=True)

nl = pill(s, "Live → gridright.netlify.app", Inches(0.5), Inches(5.2), fill=AMBER, text_color=AMBER_INK)
pill(s, "Built by Gaurav Poudel", nl, Inches(5.2), fill=PURPLE, text_color=WHITE)
footer(s, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "The Problem")
heading(s, "Your solar panels sit idle half the day.", size=38)

txt(s, "A typical rooftop produces more energy than a home can use — especially at noon.\n\nThat surplus just goes to waste.",
    Inches(0.5), Inches(2.05), Inches(5.8), Inches(1.9), size=21, color=INK)
txt(s, "A household that both produces and consumes energy is a  prosumer.  Today, prosumers have no good way to share their surplus.",
    Inches(0.5), Inches(4.15), Inches(5.8), Inches(1.3), size=16, color=MUTED, italic=True)

txt(s, "Meanwhile, your neighbor is buying grid power at full retail price.",
    Inches(6.75), Inches(2.05), Inches(6.0), Inches(1.0), size=21, color=INK)
txt(s, "Sun is free.  Power is not.\nThe gap between them is lost money for everyone.",
    Inches(6.75), Inches(3.15), Inches(6.0), Inches(1.1), size=23, bold=True, color=GREEN_DARK)
callout(s, "Net metering — how most countries credit surplus solar — pays only a fraction of the retail price. Communities deserve better.",
    Inches(6.75), Inches(4.5), Inches(6.05), Inches(1.35))
footer(s, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — HIDDEN PROBLEM
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "The Hidden Problem")
heading(s, "No easy way to share solar between neighbors.", size=36)
bullets(s, [
    "Feed-in tariffs pay a fraction of retail price — not real value.",
    "Sending surplus to your own community? Almost impossible today.",
    "Communities want to pool their DERs (Distributed Energy Resources) — rooftops, batteries, EVs — but the tools don't exist.",
    "Energy bills rise every year, even when the sun is shining on your block.",
], Inches(0.5), Inches(1.95), Inches(12.3), Inches(2.6), size=21)
callout(s, "50% of people can't put panels on their roof — renters, apartments, shaded homes.\nThey are locked out of the solar revolution entirely.",
    Inches(0.5), Inches(5.0), Inches(12.35), Inches(1.25), size=20)
footer(s, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION  (Seller / Pool / Operator)
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "Our Solution")
heading(s, "A community energy pool.", size=44)
txt(s, "One shared pool for a neighborhood — sellers contribute, an operator governs, everyone saves.",
    Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.55), size=21, color=MUTED)

cw, ch, ct = Inches(3.95), Inches(2.15), Inches(2.5)
cols = [Inches(0.5), Inches(4.66), Inches(8.82)]
data = [
    ("1", "Sellers", GREEN, "Households with solar contribute their surplus kWh to the shared community pool."),
    ("2", "Community Pool", GREEN, "The pool aggregates all contributions and dispatches energy where it's needed."),
    ("3", "Operator", GREEN, "A utility / co-op operator runs the policy layer, approves flows, and resolves exceptions."),
]
for (num, title, col, body), cx in zip(data, cols):
    card(s, cx, ct, cw, ch)
    circle_num(s, num, cx + Inches(0.25), ct + Inches(0.25))
    txt(s, title, cx + Inches(0.25), ct + Inches(0.9), cw - Inches(0.5), Inches(0.5),
        size=23, bold=True, color=GREEN_DARK)
    txt(s, body, cx + Inches(0.25), ct + Inches(1.42), cw - Inches(0.5), Inches(0.7),
        size=16, color=MUTED)

txt(s, "One neighborhood.  One pool.  Fair price.  Real savings.",
    Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.7),
    size=28, bold=True, color=GREEN_DARK)
footer(s, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (flow, fits within slide width)
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "How It Works — End-to-End Flow")
heading(s, "From rooftop to reward in five steps.", size=38)

nw, nh, nt = Inches(1.78), Inches(2.0), Inches(2.1)
aw = Inches(0.4)
nodes = [
    ("☀️", "Seller\ncontributes", "Surplus kWh logged to the pool", GREEN_SOFT, GREEN_BORD),
    ("\U0001F30A", "Community\nPool", "Aggregates DER supply & demand", GREEN_SOFT, GREEN_BORD),
    ("\U0001F916", "AI recommends", "Ranks import / export flows", AMBER_SOFT, AMBER_BORD),
    ("\U0001F464", "Operator\napproves", "Policy layer + exception queue", AMBER_SOFT, AMBER_BORD),
    ("⛓️", "Solana\nsettles", "Anchor program records on-chain", PURPLE_SOFT, PURPLE_BORD),
    ("\U0001F3C5", "cNFT badge", "Bubblegum milestone reward", PURPLE_SOFT, PURPLE_BORD),
]
total = len(nodes)*float(nw) + (len(nodes)-1)*float(aw)
cur = (float(W) - total) / 2.0
for i, (ic, lb, sb, fl, bd) in enumerate(nodes):
    flow_node(s, ic, lb, sb, Inches(cur), nt, nw, nh, fl, bd)
    cur += float(nw)
    if i < len(nodes)-1:
        arrow(s, Inches(cur), nt + Inches(0.78))
        cur += float(aw)

nl = pill(s, "Pool Dispatch", Inches(0.5), Inches(4.45), fill=AMBER, text_color=AMBER_INK, size=14)
nl = pill(s, "Operator Policy Layer", nl, Inches(4.45), fill=AMBER, text_color=AMBER_INK, size=14)
nl = pill(s, "Human-in-the-Loop (HITL)", nl, Inches(4.45), fill=PURPLE, text_color=WHITE, size=14)
pill(s, "Compressed NFT (cNFT)", nl, Inches(4.45), fill=PURPLE, text_color=WHITE, size=14)
footer(s, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — AI + HUMAN IN THE LOOP
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "AI Architecture — The Human-in-the-Loop Principle")
heading(s, "AI recommends.  Operator decides.", size=40)

card(s, Inches(0.5), Inches(1.9), Inches(5.95), Inches(3.15), fill=PURPLE_SOFT, border=PURPLE_BORD)
txt(s, "\U0001F916  What the AI engine does", Inches(0.72), Inches(2.05), Inches(5.6), Inches(0.5),
    size=20, bold=True, color=PURPLE)
bullets(s, [
    "Reads real-time pool state — supply, demand, price signals",
    "Forecasts solar production curves over the day",
    "Ranks import / export flows by cost and pool balance",
    "Flags anomalies to the operator's exception queue",
], Inches(0.72), Inches(2.55), Inches(5.55), Inches(2.3), size=17, gap=Pt(6))

card(s, Inches(6.9), Inches(1.9), Inches(5.95), Inches(3.15))
txt(s, "\U0001F464  What the operator does", Inches(7.12), Inches(2.05), Inches(5.6), Inches(0.5),
    size=20, bold=True, color=GREEN_DARK)
bullets(s, [
    "Reviews the AI's ranked recommendations — never a done deal",
    "Runs the policy layer + exception queue: disputes, anomalies",
    "Approves or rejects each flow — the final call is always human",
    "Sets pool policy: price floor, max draw per household",
], Inches(7.12), Inches(2.55), Inches(5.55), Inches(2.3), size=17, gap=Pt(6))

callout(s, "Why it matters: autonomous AI settlement is how you make the front page for the wrong reason. Trust is the moat — and trust requires a human in the loop.",
    Inches(0.5), Inches(5.3), Inches(12.35), Inches(1.05), size=19)
footer(s, 6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT WE BUILT
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "What We Built")
heading(s, "Three things, working together.", size=42)
steps = [
    ("A", "Pool App", "Next.js + Supabase",
     "Separate seller and operator dashboards. Two-layer auth (seller JWT + operator-only), with Row-Level Security (RLS) so each user sees only their own data."),
    ("B", "Smart Engine", "Python FastAPI + AI",
     "AI recommends import / export flows across the pool. The operator policy layer + exception queue reviews every recommendation — HITL by design. AI never settles."),
    ("C", "Settlement Layer", "Anchor + Metaplex Bubblegum",
     "Every settlement is recorded on Solana via a custom Anchor program. Sellers earn cNFT badges at contribution milestones — 100 kWh earns the first."),
]
top = Inches(1.9)
for letter, title, sub, body in steps:
    circle_num(s, letter, Inches(0.5), top + Inches(0.05), d=Inches(0.62), size=22)
    txt(s, f"{title}   —   {sub}", Inches(1.32), top, Inches(11.4), Inches(0.45),
        size=21, bold=True, color=INK)
    txt(s, body, Inches(1.32), top + Inches(0.46), Inches(11.4), Inches(0.9),
        size=17, color=MUTED)
    top += Inches(1.55)
footer(s, 7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "How We Built It")
heading(s, "Boring tools.  Interesting idea.", size=42)
txt(s, "We picked reliable, mainstream tools — so the only hard part is the idea.",
    Inches(0.5), Inches(1.78), Inches(12), Inches(0.5), size=20, color=MUTED)

rows = [
    ("Web app",         "Next.js (React)",            "Seller + operator dashboards, deployed to Netlify"),
    ("Backend API",     "Python / FastAPI",           "AI forecasting + pool policy logic, hosted on Render"),
    ("Database & auth", "Supabase (Postgres)",        "Two-layer auth + RLS per-user isolation"),
    ("Blockchain",      "Solana / Anchor",            "Custom program — every settlement on-chain"),
    ("cNFT badges",     "Metaplex Bubblegum",         "Compressed NFTs — cheap to mint, on-chain proof"),
    ("Scheduler",       "cron-job.org",               "Hourly forecasts, daily commitments, keep-alive"),
    ("Hosting",         "Netlify · Render · Supabase", "100% free tier — no credit card"),
]
x0, x1, x2 = Inches(0.6), Inches(3.1), Inches(6.2)
ht = Inches(2.35)
txt(s, "LAYER",      x0, ht, Inches(2.4), Inches(0.3), size=12, bold=True, color=MUTED)
txt(s, "TECHNOLOGY", x1, ht, Inches(3.0), Inches(0.3), size=12, bold=True, color=MUTED)
txt(s, "KEY DETAIL", x2, ht, Inches(6.4), Inches(0.3), size=12, bold=True, color=MUTED)
ry = ht + Inches(0.36)
for i, (a, b, c) in enumerate(rows):
    rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.5),
         fill=(SLATE_SOFT if i % 2 == 0 else WHITE))
    txt(s, a, x0, ry + Inches(0.07), Inches(2.4), Inches(0.4), size=16, bold=True, color=INK)
    txt(s, b, x1, ry + Inches(0.07), Inches(3.0), Inches(0.4), size=16, bold=True, color=GREEN_DARK)
    txt(s, c, x2, ry + Inches(0.07), Inches(6.5), Inches(0.4), size=16, color=MUTED)
    ry += Inches(0.5)
txt(s, "Fully deployed. No credit card. No mocks. Real users could sign up today.",
    Inches(0.5), ry + Inches(0.12), Inches(12), Inches(0.4),
    size=16, color=MUTED, italic=True)
footer(s, 8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — MARKET
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "The Market")
heading(s, "Community solar is exploding.", size=44)

def stat(sx, sy, sw, big, label):
    card(s, sx, sy, sw, Inches(1.25))
    txt(s, big, sx + Inches(0.2), sy + Inches(0.12), sw - Inches(0.4), Inches(0.6),
        size=30, bold=True, color=GREEN_DARK)
    txt(s, label, sx + Inches(0.2), sy + Inches(0.74), sw - Inches(0.4), Inches(0.45),
        size=13, color=MUTED)

r1 = Inches(1.95)
stat(Inches(0.5),  r1, Inches(3.9), "$3.2B", "US community solar market, 2025")
stat(Inches(4.62), r1, Inches(3.9), "19%",   "Annual growth, year-over-year")
stat(Inches(8.74), r1, Inches(4.1), "1.3M",  "US households already served")
r2 = Inches(3.4)
stat(Inches(0.5),  r2, Inches(6.0), "5M",          "US homes targeted by 2030 (EPA Solar for All)")
stat(Inches(6.72), r2, Inches(6.11), "$7B → $24B", "Global Virtual Power Plant market by 2030")

callout(s, "A Virtual Power Plant (VPP) is a network of distributed energy resources acting as one coordinated source. GridRight is the community-scale VPP layer that doesn't yet exist.",
    Inches(0.5), Inches(5.05), Inches(12.35), Inches(1.1), size=19)
txt(s, "Sources: NREL · Wood Mackenzie · EPA Solar for All · BloombergNEF 2025 · Rocky Mountain Institute.",
    Inches(0.5), Inches(6.9), Inches(12), Inches(0.35), size=12, color=MUTED)
footer(s, 9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — WHERE THIS WORKS TODAY (no flag emojis)
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "Where This Works Today")
heading(s, "18 countries already pay for surplus solar.", size=36)
txt(s, "Net metering, feed-in tariffs, export guarantees — live right now. GridRight adds the sharing layer on top.",
    Inches(0.5), Inches(1.72), Inches(12.4), Inches(0.5), size=18, color=MUTED)

left_c = [
    ("USA", "Net metering (state-by-state)"),
    ("Australia", "Feed-in tariff"),
    ("Germany", "EEG feed-in tariff"),
    ("Japan", "FIT surplus purchase"),
    ("United Kingdom", "Smart Export Guarantee"),
    ("Italy", "Scambio Sul Posto"),
    ("Spain", "Surplus compensation"),
    ("France", "Autoconsommation sale"),
    ("Belgium", "Flanders net metering"),
]
right_c = [
    ("Denmark", "Net metering up to 6 kW"),
    ("Canada", "Ontario, Alberta, BC"),
    ("India", "Rooftop net metering"),
    ("Brazil", "REN 482/2012"),
    ("Mexico", "Residential net metering"),
    ("South Africa", "Feed-in tariff"),
    ("Israel", "Net metering"),
    ("Netherlands", "Saldering net metering"),
    ("South Korea", "RPS surplus purchase"),
]

def country_col(items, lx):
    ty = Inches(2.35)
    for name, prog in items:
        rect(s, lx, ty + Inches(0.11), Inches(0.13), Inches(0.13), fill=GREEN)  # bullet dot
        tb = s.shapes.add_textbox(lx + Inches(0.28), ty, Inches(5.8), Inches(0.4))
        tf = tb.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = name + "  "; _style_run(r1, 17, True, INK)
        r2 = p.add_run(); r2.text = "—  " + prog; _style_run(r2, 17, False, MUTED)
        ty += Inches(0.44)

country_col(left_c,  Inches(0.5))
country_col(right_c, Inches(6.85))
txt(s, "Sources: IEA Renewables 2024 · Fraunhofer ISE · Clean Energy Council AU · METI Japan · Ofgem UK.",
    Inches(0.5), Inches(6.95), Inches(12), Inches(0.35), size=12, color=MUTED)
footer(s, 10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — FUTURE & BUSINESS MODEL
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "Future & Business Model")
heading(s, "Where this goes next.", size=44)

card(s, Inches(0.5), Inches(1.95), Inches(5.95), Inches(3.45))
txt(s, "\U0001F680  Next steps", Inches(0.72), Inches(2.08), Inches(5.6), Inches(0.5),
    size=22, bold=True, color=GREEN_DARK)
bullets(s, [
    "Pilot with a US community solar program (MN, NY)",
    "Add battery storage + EV charging to the pool",
    "White-label for utilities and energy co-ops",
    "VPP interoperability via the OpenADR standard",
], Inches(0.72), Inches(2.58), Inches(5.5), Inches(2.6), size=17, gap=Pt(6))

card(s, Inches(6.9), Inches(1.95), Inches(5.95), Inches(3.45))
txt(s, "\U0001F4B0  How we make money", Inches(7.12), Inches(2.08), Inches(5.6), Inches(0.5),
    size=22, bold=True, color=GREEN_DARK)
bullets(s, [
    "1–2% transaction fee on every pool settlement",
    "SaaS subscription for utility / co-op operators",
    "Premium AI forecasting for grid operators",
    "cNFT badge marketplace & sponsorships",
], Inches(7.12), Inches(2.58), Inches(5.5), Inches(2.6), size=17, gap=Pt(6))

txt(s, "Small fees  ·  Large volume  ·  Real network effects.",
    Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.6),
    size=26, bold=True, color=GREEN_DARK)
footer(s, 11)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — WHY THIS WINS
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s, WHITE); left_bar(s)
eyebrow(s, "Why This Wins")
heading(s, "Rare.  Proven.  Live.  Inclusive.", size=44)
bullets(s, [
    "Rare — almost no one runs an operator-governed community solar pool with real on-chain settlement.",
    "Proven market — 18 countries already pay for surplus solar; we add the sharing layer on top.",
    "Live today — not a slide deck, a running app.  gridright.netlify.app  →  sign up now.",
    "Inclusive — works for renters and apartments, not just homeowners. Anyone in the pool benefits.",
    "Open by design — HITL: AI recommends, the operator decides. Trust is the moat scale can't copy.",
    "First-mover — community solar + Solana settlement + cNFT badges in a single, deployed product.",
], Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.9), size=18, gap=Pt(5))
txt(s, "Thank you.", Inches(0.5), Inches(5.95), Inches(7), Inches(0.85),
    size=46, bold=True, color=GREEN_DARK)
txt(s, "gridright.netlify.app  ·  github.com/Gauravpoudel7/gridright  ·  Built by Gaurav Poudel",
    Inches(0.5), Inches(6.8), Inches(12.2), Inches(0.45), size=15, color=MUTED, italic=True)
footer(s, 12)


# ── Save ───────────────────────────────────────────────────────
OUT = r"C:\Users\gaura\OneDrive\Desktop\gridright\hackathon-presentation\GridRight.pptx"
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
